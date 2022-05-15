'''
Replicate "Defining seizure onsets and endings" seciton of "schindler 2006 Brain

'''
#%%
import os
import sys
import json
import numpy as np
from scipy.io import loadmat
from os.path import join as ospj
import pandas as pd
import matplotlib.pyplot as plt
from tqdm import tqdm
import seaborn as sns
from pptx import Presentation
from scipy.stats import iqr
# sns.set_style("whitegrid")
sns.set_context('notebook')
from scipy.signal import decimate

import tools

SHAT_THRESHOLD = 3
WINDOW_THRESHOLD = 0.1

# %%
# # Get paths from config file and metadata

code_path = os.path.dirname(__file__)
with open(ospj(code_path, "config.json"), 'rb') as f:
# with open("config.json", 'rb') as f:
    config = json.load(f)
repo_path = config['repositoryPath']
data_path = ospj(repo_path, 'data')
figure_path = ospj(repo_path, 'figures')
electrodes_opt = config['electrodes']

metadata_path = config['metadataPath']
atlas_path = config['atlas_path']

# load in seizure metadata
sz_metadata = pd.read_excel(ospj(data_path, "seizure_metadata_all_atlas.xlsx"))
# sz_metadata = sz_metadata[sz_metadata['Seizure duration'] < 300]
# remove HUP170 maybe?
sz_metadata = sz_metadata[sz_metadata.Patient != 'HUP170']

patient_cohort = pd.read_excel(ospj(data_path, 'patient_cohort_thresholds.xlsx'))
# for making figures
wd_ratio = 0.0005
ht_ratio = 0.2

# %%
def moving_average(x, w):
    '''
    This function calculates the average of x over w windows.
    '''
    x_non_nan = np.nan_to_num(x)
    (n_time_windows, n_channels) = x_non_nan.shape
    out_data = np.zeros((n_time_windows - (w - 1), n_channels))
    for i in range(x.shape[-1]):
        out_data[:, i] = np.convolve(x_non_nan[:, i], np.ones(w), 'valid') / w
    return out_data

nan_indices = []
# %% Determine recruited and unrecruited channels
for index, row in sz_metadata.iterrows():
    pt = row['Patient']
    sz_num = row['Seizure number']
    sz_category = row['Seizure category']

    pt_threshold = patient_cohort[patient_cohort['Patient'] == pt]['threshold'].values[0]

    # pt = "HUP157"
    # sz_num = 1
    # sz_category = "Focal"

    print(pt, sz_num)
    pt_data_path = ospj(data_path, pt)
    pt_figure_path = ospj(figure_path, pt)

    # CHANGE TO NOT
    if not os.path.exists(ospj(pt_data_path, f"Shat_matrix_sz-{sz_num}_period-ictal.pkl")):
        # Calculate S(t) for preictal clips
        fname = f"raw_signal_elec-{electrodes_opt}_period-preictal_sz-{sz_num}.pkl"
        pi_df = pd.read_pickle(ospj(data_path, pt, fname))
        (n_time_windows, n_channels) = pi_df.shape
        time_pi = pi_df.index.total_seconds()
        Fs = round(1 / (time_pi[1] - time_pi[0]))

        # Find windows where there are artifacts
        win_size = int(0.5 * Fs)
        ind_overlap = np.reshape(np.arange(pi_df.shape[0]), (-1, int(win_size)))
        n_windows = np.size(ind_overlap, axis=0)

        is_nan = np.ones((n_time_windows, n_channels), dtype=bool)
        for win_inds in ind_overlap:
            is_nan[win_inds, :] = np.sum(np.isnan(pi_df.iloc[win_inds, :]))
            is_nan[win_inds, :] = (np.sum(np.abs(pi_df.iloc[win_inds, :]), axis=0) < 1/12)
            is_nan[win_inds, :] = (np.sqrt(np.sum(np.power(np.diff(pi_df.iloc[win_inds, :], axis=0), 2), axis=0)) > 15000)

        # Calculate absolute slope
        S_pi = np.abs(np.divide(np.diff(pi_df, axis=0), 1/Fs))

        if np.sum(~is_nan) != 0:
            nan_indices.append(index)

        # calculate std on non-artifact windows and channels
        S_pi[is_nan[1:, :]] = np.NaN
        sigma_pi = np.nanstd(S_pi, axis=0)
        sigma_pi = np.mean(sigma_pi)
        # sigma_pi = iqr(S_pi, axis=0, nan_policy='omit')

        # Calculate S(t) for ictal clips
        fname = f"raw_signal_elec-{electrodes_opt}_period-ictal_sz-{sz_num}.pkl"
        sz_df = pd.read_pickle(ospj(data_path, pt, fname))
        (n_time_windows, n_channels) = sz_df.shape
        time_ic = sz_df.index.total_seconds()
        Fs = round(1 / (time_ic[1] - time_ic[0]))

        # Find windows where there are artifacts
        win_size = int(0.5 * Fs)
        ind_overlap = np.reshape(np.arange(len(time_ic)), (-1, int(win_size)))
        n_windows = np.size(ind_overlap, axis=0)

        is_nan = np.ones((n_time_windows, n_channels), dtype=bool)
        for win_inds in ind_overlap:
            is_nan[win_inds, :] = np.sum(np.isnan(sz_df.iloc[win_inds, :])).any()
            is_nan[win_inds, :] = (np.sum(np.abs(sz_df.iloc[win_inds, :]), axis=0) < 1/12).any()
            is_nan[win_inds, :] = (np.sqrt(np.sum(np.power(np.diff(sz_df.iloc[win_inds, :], axis=0), 2), axis=0)) > 15000).any()

        S_ic = np.abs(np.divide(np.diff(sz_df, axis=0), 1/Fs))

        print(np.sum(is_nan))

        Shat_ic = np.divide(S_ic, sigma_pi)

        # smooth by 5% of seizure duration
        smooth_window = int(Fs * (time_ic[-1] - time_ic[0]) / 20)
        Shat_ic = moving_average(Shat_ic, smooth_window)
        time_ic = time_ic[(smooth_window):]

        # Shat_ic[is_nan[1:, :]] = np.NaN

        Shat_ic = pd.DataFrame(data=Shat_ic, index=time_ic[:], columns=sz_df.columns).dropna()
        Shat_ic.to_pickle(ospj(pt_data_path, f"Shat_matrix_sz-{sz_num}_period-ictal.pkl"))

    else:
        Shat_ic = pd.read_pickle(ospj(pt_data_path, f"Shat_matrix_sz-{sz_num}_period-ictal.pkl"))

    recruited = np.sum(Shat_ic > pt_threshold) > Shat_ic.shape[0] * WINDOW_THRESHOLD

    recruited.to_csv(ospj(pt_data_path, f"recruited_schindler_shatthresh-{SHAT_THRESHOLD}_winthresh-{WINDOW_THRESHOLD}_sz-{sz_num}.csv"))
    np.save(ospj(pt_data_path, f"recruited_schindler_shatthresh-{SHAT_THRESHOLD}_winthresh-{WINDOW_THRESHOLD}_sz-{sz_num}.npy"), np.array(recruited))

    # break
# %% Determine recruited and unrecruited channels
for index, row in sz_metadata.iterrows():
    pt = row['Patient']
    sz_num = row['Seizure number']
    sz_category = row['Seizure category']
    sz_eec = row['Seizure EEC']
    sz_ueo = row['Seizure UEO']

    # if sz_num != 1:
    #     continue
    print(pt, sz_num)

    # pt = "HUP151"
    # sz_num = 9
    # sz_category = "Focal"

    pt_data_path = ospj(data_path, pt)
    pt_figure_path = ospj(figure_path, pt)

    fname = "raw_signal_elec-{}_period-ictal_sz-{}.pkl".format(electrodes_opt, sz_num)

    sz_df = pd.read_pickle(ospj(data_path, pt, fname))
    (n_time_windows, n_channels) = sz_df.shape

    sz_time = sz_df.index.total_seconds()

    Fs = round(1 / (sz_time[1] - sz_time[0]))

    fname = "raw_signal_elec-{}_period-preictal_sz-{}.pkl".format(electrodes_opt, sz_num)
    preictal_df = pd.read_pickle(ospj(data_path, pt, fname))
    (n_time_windows, n_channels) = preictal_df.shape

    preictal_time = preictal_df.index.total_seconds()

    total = pd.concat([preictal_df, sz_df])
    time = np.concatenate([preictal_time, sz_time])

    total = decimate(total, 10, axis=0)
    time = np.linspace(time[0], time[-1], total.shape[0])

    recruited  = np.load(ospj(pt_data_path, f"recruited_schindler_shatthresh-{SHAT_THRESHOLD}_winthresh-{WINDOW_THRESHOLD}_sz-{sz_num}.npy"))

    fig, ax = plt.subplots(figsize=(n_time_windows*wd_ratio,n_channels*ht_ratio))

    recruited_inds = np.where(recruited)[0]
    nonrecruited_inds = np.where(~recruited)[0]

    if len(recruited_inds) != 0:
        ax.plot(time, total[:, recruited] + 1000*recruited_inds, color='r')
    if len(nonrecruited_inds) != 0:
        ax.plot(time, total[:, ~recruited] + 1000*nonrecruited_inds, color='k')

    # ax.plot(time, total + 1000*np.arange(n_channels), color=colors
    # )
    ax.set_title("{}, {} ({})".format(pt, sz_num, sz_category))
    ax.set_xlabel('EMU Time (s)')
    ax.set_yticks([])    # ax.axvline(sz_time[0])

    ax.axvline(sz_eec)
    # ax.axvline(sz_ueo)
    fname = "raw_signal_recruited_schindler_labels_elec-{}_period-ictal_sz-{}".format(electrodes_opt, sz_num)
    plt.savefig(ospj(pt_figure_path, "{}.svg".format(fname)), bbox_inches='tight', transparent='true')
    plt.savefig(ospj(pt_figure_path, "{}.png".format(fname)), bbox_inches='tight', transparent='true')
    
    plt.close()
#%%
prs = Presentation()
for index, row in sz_metadata.iterrows():
    pt = row['Patient']
    sz_num = row['Seizure number']

    # pt = "HUP151"
    # sz_num = 11
    # sz_category = "Focal"

    pt_figure_path = ospj(figure_path, pt)
    fname = "raw_signal_recruited_schindler_labels_elec-{}_period-ictal_sz-{}".format(electrodes_opt, sz_num)

    if not os.path.exists(ospj(pt_figure_path, '{}.png'.format(fname))):
        continue
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    tools.add_img_to_ppt(slide, 1, ospj(pt_figure_path, '{}.png'.format(fname)))
    title_placeholder = slide.shapes.title
    title_placeholder.text = "{}, Seizure {}".format(pt, sz_num)    
    
prs.save(ospj(figure_path, 'raw_signals_recruited_labels_period-ictal.pptx'))

# %%
